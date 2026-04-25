"""Build STAR final presentation by populating the unpacked-then-repacked
14-slide skeleton with content, images, and speaker notes.

Input:  /tmp/deck_skeleton.pptx (skeleton with 14 slides in correct order)
Output: final_docs/final_presentation/STAR_Final_Presentation.pptx
"""

import copy
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Emu, Pt
from pptx.dml.color import RGBColor

REPO = Path("/Users/cesarmagana/Documents/GitHub/STAR")
SKELETON = "/tmp/deck_skeleton.pptx"
OUT = REPO / "final_docs/final_presentation/STAR_Final_Presentation.pptx"

CHARTS = REPO / "final_capstone/charts/output"
SCREENS = REPO / "final_docs/system_software_user_manual/images"

# === Per-slide content =====================================================

SLIDES = [
    # 1: Title (Title Slide layout)
    {
        "kind": "title",
        "date": "Apr 28, 2026",
        "title_top": "STAR",
        "title_main": "Spatial Topography\nAccessibility Robot",
        "presenters": "Locked Inc.  ESET 420  Dr. Logan Porter",
        "footer": "Texas A&M University\nESET Senior Capstone, Spring 2026",
        "notes": (
            "Hello, we are Locked Inc. STAR is the Spatial Topography "
            "Accessibility Robot -- an autonomous indoor robot that "
            "audits buildings for ADA Title III compliance. Today we "
            "will walk you through the problem, our solution, the "
            "system architecture, validation results, and a live demo."
        ),
    },
    # 2: $236,451 hook (Two-content)
    {
        "kind": "two_col",
        "title": "$236,451 per door",
        "left": [
            "Federal civil penalty per repeat ADA violation, after the July 2025 Federal Register inflation adjustment.",
            "8,800 federal Title III lawsuits were filed in 2024 (Seyfarth Shaw tracker), continuing a sustained 9k-11k baseline since 2018.",
            "Manual audits cost $800-$50,000+ per facility and take days. The math does not close.",
        ],
        "image": None,
        "image_caption": "$118,225\nfirst violation\n\n$236,451\neach subsequent",
        "image_caption_style": "stat_callout",
        "notes": (
            "The federal civil penalty for a repeat ADA Title III "
            "violation is now $236,451 per occurrence as of the July "
            "2025 Federal Register adjustment. In 2024 there were "
            "8,800 federal Title III lawsuits filed. Manual audits "
            "by certified inspectors cost between $800 for a small "
            "store to over $50,000 for a campus, and take two to "
            "five days each. The economics push everything toward "
            "automation."
        ),
    },
    # 3: Who measures? (Two-content with donut chart)
    {
        "kind": "two_col",
        "title": "Who measures this?",
        "left": [
            "73 percent of public-facing facilities fail at least one ADA standard at survey.",
            "31 million U.S. adults (12.2 percent) report mobility disabilities, per CDC.",
            "Only 32 percent of U.S. property buyers ask about ADA compliance during diligence.",
        ],
        "image": str(CHARTS / "chart_b_compliance_donut.png"),
        "image_caption": None,
        "notes": (
            "Despite decades of legislation, 73% of public-facing "
            "facilities fail at least one ADA standard. 31 million "
            "Americans have mobility disabilities. Only 32% of "
            "property buyers ask about ADA at all. The audit gap is "
            "real and it is structural."
        ),
    },
    # 4: Enforcement accelerating (Two-content with lawsuit chart)
    {
        "kind": "two_col",
        "title": "Enforcement is accelerating",
        "left": [
            "Federal Title III filings: 10,163 (2018) -> peak 11,452 (2021) -> rebounded to 8,800 in 2024.",
            "California leads with 3,252 filings in 2024; Texas filed 224.",
            "Title II web-accessibility deadline (April 2026 for state and local) drove a 7 percent year-over-year increase.",
        ],
        "image": str(CHARTS / "chart_a_lawsuit_trend.png"),
        "image_caption": None,
        "notes": (
            "Filings have stayed in the 8,000 to 11,500 per year band "
            "for seven years now. They peaked at 11,452 in 2021, fell, "
            "and rebounded to 8,800 in 2024 - a 7% jump - driven in "
            "part by the new Title II web accessibility deadlines for "
            "state and local entities. California is the clear leader "
            "with 3,252 filings; Texas had 224. Enforcement is not "
            "fading, it is becoming routine."
        ),
    },
    # 5: Tech stack (3-column)
    {
        "kind": "three_col",
        "title": "Tech stack",
        "cols": [
            {
                "head": "Hardware",
                "body": [
                    "RPLiDAR C1 (12 m, IP54)",
                    "IMX219-83 stereo cameras",
                    "BNO055 IMU, HC-SR04 sonar",
                    "Renesas RX72N (240 MHz)",
                    "Raspberry Pi 5",
                    "DRV8263H driver x4",
                    "5S 18650 battery pack",
                ],
            },
            {
                "head": "Software",
                "body": [
                    "ROS 2 Jazzy on Ubuntu 24.04",
                    "slam_toolbox (async)",
                    "Nav2 + DWB controller",
                    "robot_localization EKF",
                    "Go gateway with nanopb",
                    "Python compliance engine",
                    "C23 firmware on ThreadX",
                ],
            },
            {
                "head": "Validation",
                "body": [
                    "60 firmware ctests, 100 pct libs/",
                    "74.0 pct Go gateway cov, gate 65",
                    "14 ROS 2 colcon tests",
                    "12 compliance pytest modules",
                    "9 manual PCB validation checks",
                    "Wixey WR300 +/- 0.1 deg ground truth",
                ],
            },
        ],
        "notes": (
            "STAR is a multi-language stack. Hardware: a custom 6-layer "
            "PCB built around the Renesas RX72N for real-time motor "
            "control, a Pi 5 for autonomy, RPLiDAR C1 for 2D scan, "
            "stereo cameras for 3D, BNO055 IMU. Software: ROS 2 Jazzy "
            "for autonomy, slam_toolbox plus Nav2 for navigation, a "
            "Go gateway for telemetry, a Python compliance engine. "
            "Validation: 60 firmware Unity tests at 100% library "
            "coverage, 74% Go coverage gated at 65%, 14 ROS 2 tests, "
            "12 pytest modules, plus formal manual PCB validation."
        ),
    },
    # 6: System architecture (Two-content with sankey)
    {
        "kind": "two_col",
        "title": "System architecture",
        "left": [
            "Sensors feed the RX72N for hard-real-time motor control at 250 Hz.",
            "RX72N talks to the Pi 5 over a framed transport: SPI 10 MHz with FEC and HARQ is the architected production link; USB-CDC ASCII at 115200 baud is the deployed link for capstone defense.",
            "Pi 5 runs slam_toolbox, robot_localization EKF, and Nav2. The compliance engine reads the fused map and emits a structured PDF report.",
        ],
        "image": str(CHARTS / "chart_f_pipeline_sankey.png"),
        "image_caption": None,
        "notes": (
            "The pipeline: sensors stream into the RX72N which closes "
            "motor loops at 250 Hz. The RX72N talks to the Pi 5 via a "
            "framed serial link - architected as SPI at 10 MHz with "
            "FEC and HARQ for production, but deployed for defense as "
            "a simpler USB-CDC ASCII bridge at 115,200 baud. The Pi 5 "
            "runs slam_toolbox, an EKF that fuses wheel odometry with "
            "the IMU, and Nav2 for path planning. The compliance "
            "engine sits on top, reads the fused map, and emits a PDF "
            "audit report."
        ),
    },
    # 7: 7 ADA checks (3-column)
    {
        "kind": "three_col",
        "title": "Seven ADA checks",
        "cols": [
            {
                "head": "Implemented",
                "body": [
                    "Ramp slope (ADA 405.2)",
                    "RANSAC plane on LiDAR + BNO055 pitch cross-validation",
                    "Threshold 4.76 deg (1:12)",
                ],
            },
            {
                "head": "Stretch",
                "body": [
                    "Trip hazard (ADA 303): vertical discontinuity > 1/4 in",
                    "Path width (ADA 403.5): cross-section along planned route",
                ],
            },
            {
                "head": "Architected",
                "body": [
                    "Ramp width (405.5)",
                    "Ramp landings (405.7)",
                    "Door clear width (404.2.3)",
                    "Door threshold (404.2.5)",
                    "Node skeletons committed; bodies pending",
                ],
            },
        ],
        "notes": (
            "Seven ADA geometric checks, mapped to three status tiers. "
            "Ramp slope is fully IMPLEMENTED with measured ground "
            "truth. Trip hazard and path width are STRETCH - "
            "scaffolding exists, full algorithm pending. The other "
            "four are ARCHITECTED: every node has a Python file in "
            "the repo, message schemas exist, but the geometric logic "
            "is empty."
        ),
    },
    # 8: Demo (Two-content with lichtblick)
    {
        "kind": "two_col",
        "title": "Live demo",
        "left": [
            "Robot is dispatched in MANUAL via the Grafana cockpit, then handed to AUTONOMY for frontier exploration.",
            "On approach to a ramp, the compliance engine fits a RANSAC plane and cross-validates pitch against the BNO055.",
            "When measured slope > 4.76 deg, the engine logs a violation flag with the ground-truth row in the PDF audit report.",
        ],
        "image": str(SCREENS / "lichtblick-teleop.png"),
        "image_caption": None,
        "notes": (
            "For the live demo: we start the robot in MANUAL mode "
            "from the Grafana cockpit. We hand it off to AUTONOMY, "
            "Nav2 takes over, and frontier exploration drives the "
            "robot into a hallway. As it approaches the ramp, the "
            "compliance engine fits a plane to the LiDAR returns "
            "using RANSAC and cross-validates that against the "
            "BNO055 IMU pitch. If the agreement is within "
            "0.5 degrees and the slope exceeds 4.76 degrees, the "
            "engine flags a Title III violation and logs the "
            "ground-truth measurement to the audit report."
        ),
    },
    # 9: Validation (Two-content with thompson hall)
    {
        "kind": "two_col",
        "title": "Validation",
        "left": [
            "Software: 60 firmware ctests (100 pct libs/), 74 pct Go gateway, 14 ROS 2 tests, 12 pytest modules.",
            "Hardware: 9 manual PCB checks all PASS -- continuity, motor free-spin + loaded chassis with DC current, scoped PWM, AD2-probed comm.",
            "Field: 200 sq ft autonomous scan in Thompson Hall (right). Two issues detected: slope > 1:12, path < 36 in.",
        ],
        "image": str(SCREENS / "thompson-hall-map.png"),
        "image_caption": "Thompson Hall, 200 sq ft.",
        "notes": (
            "Validation in three layers. Software: 60 firmware Unity "
            "tests passing at 100% library coverage, 74% gateway "
            "coverage, 14 ROS 2 tests, 12 compliance pytest modules. "
            "Hardware: nine manual PCB validation checks, all passing "
            "first-pass, including motor free-spin and loaded-chassis "
            "with current measurements and scope-probed PWM. Field: "
            "we mapped 200 square feet of Thompson Hall on the Texas "
            "A&M campus and detected two ADA compliance issues."
        ),
    },
    # 10: Cost (Two-content with cost-time chart)
    {
        "kind": "two_col",
        "title": "Cost",
        "left": [
            "Total BOM: ~$864 with PCB run, well under the $2,000 capstone target.",
            "Compute $182 (Pi 5, RX72N, PCB) - Sensors $270 (LiDAR, stereo, IMU, sonar) - Drivetrain $144 (motors, drivers, wheels) - Power $108 (battery, BMS, regulators) - Chassis $100",
            "Daxbot service-as-a-product is $920k for 950 surveyed miles. STAR is sub-$2k recurring near-zero. Different economics entirely.",
        ],
        "image": str(CHARTS / "chart_c_cost_time.png"),
        "image_caption": None,
        "notes": (
            "Total BOM is approximately $864 with the PCB run "
            "included - well under our $2000 capstone target. The "
            "biggest line is sensors at $270, then compute at $182, "
            "drivetrain at $144, power and chassis at about $100 "
            "each. For comparison, Daxbot, the closest commercial "
            "service, is $920,000 for 950 miles of surveying. STAR "
            "fits a completely different economic envelope."
        ),
    },
    # 11: Section breaker - Broader impact
    {
        "kind": "section",
        "title": "Broader impact",
        "notes": (
            "Now let us talk about who STAR helps and where it could "
            "go from here."
        ),
    },
    # 12: Future work (Title + bullets)
    {
        "kind": "bullets",
        "title": "Future work",
        "bullets": [
            "Implement architected ADA checks: ramp width, ramp landing, door clear width, door threshold.",
            "Sensor fusion via RTAB-Map for richer 3D reconstruction.",
            "Multi-floor support with elevator detection and floor handoff.",
            "ML-based fixture recognition for fire extinguishers, signage, hardware.",
            "Outdoor GPS plus RTK for campus-scale audit campaigns.",
        ],
        "notes": (
            "Future work breaks down into four buckets. First, finish "
            "the architected ADA checks - the schemas exist, we just "
            "need the geometric algorithms. Second, fuse RGB-D and "
            "LiDAR via RTAB-Map for richer 3D reconstruction. Third, "
            "multi-floor support with elevator detection. Fourth, "
            "ML-driven fixture recognition for things like fire "
            "extinguishers and signage. And eventually, outdoor "
            "GPS+RTK so we can audit a whole campus, not just a "
            "single floor."
        ),
    },
    # 13: Q&A team (3-column)
    {
        "kind": "three_col",
        "title": "Q&A and team",
        "cols": [
            {
                "head": "Software",
                "body": [
                    "Cesar Magana, Project Manager",
                    "proto, gateway, ROS 2, UI, MATLAB",
                ],
            },
            {
                "head": "Embedded and PCB",
                "body": [
                    "Brighton Sikarskie",
                    "firmware, compliance engine, infra, schematic",
                    "Jeremie Hockey, PCB layout",
                ],
            },
            {
                "head": "Hardware and Review",
                "body": [
                    "Michael Norton, mechanical design",
                    "Shawn Ciaciura, schematic review",
                    "Jared Bartz, schematic review",
                ],
            },
        ],
        "notes": (
            "We are happy to take questions. The team is six people. "
            "I am Cesar Magana, project manager and software lead. "
            "Brighton Sikarskie ran firmware, compliance engine, "
            "infra, and schematic. Jeremie Hockey did PCB layout. "
            "Michael Norton owned mechanical design. Shawn Ciaciura "
            "and Jared Bartz contributed schematic review. Source "
            "code is open at github.com/Locked-Inc/STAR."
        ),
    },
    # 14: Thank You (Two Content closing)
    {
        "kind": "thanks",
        "footer": "Locked Inc.  -  ESET 420 Senior Capstone, Spring 2026  -  github.com/Locked-Inc/STAR",
        "notes": "",
    },
]

# === Helpers ================================================================


def replace_paragraph_text(paragraph, new_text):
    """Replace paragraph text while keeping the formatting of the first run.

    Removes any extra runs after the first, then sets the first run's text.
    """
    runs = paragraph.runs
    if not runs:
        # Empty paragraph; create a run.
        run = paragraph.add_run()
        run.text = new_text
        return run
    first = runs[0]
    # Drop runs after the first by removing their <a:r> from the XML.
    for extra_run in runs[1:]:
        extra_run._r.getparent().remove(extra_run._r)
    first.text = new_text
    return first


def set_text_block(text_frame, lines, *, first_bold=False, first_uppercase=False,
                   bullet_glyph=None):
    """Replace the text frame's content with N paragraphs (one per line).

    Preserves the first paragraph's first-run formatting (size, color, font),
    duplicates that paragraph for each subsequent line, then sets each line's
    text. If `first_bold` is True, only the first paragraph is bolded
    (useful for SUBHEADING + body in template's three-column layout).
    """
    paragraphs = text_frame.paragraphs
    if not paragraphs:
        return
    # Preserve a deep copy of the first paragraph's XML so we can clone it
    # for each new line. This keeps run formatting intact.
    template_para = paragraphs[0]
    template_xml = copy.deepcopy(template_para._p)
    # Remove all existing paragraphs except the first.
    body_elt = template_para._p.getparent()
    for p in paragraphs[1:]:
        body_elt.remove(p._p)
    # Now write line[0] into the existing paragraph, and append clones for
    # subsequent lines.
    for i, line in enumerate(lines):
        if i == 0:
            target_p = template_para._p
        else:
            new_p = copy.deepcopy(template_xml)
            body_elt.append(new_p)
            target_p = new_p
        # Wrap in a paragraph object so .runs works.
        from pptx.text.text import _Paragraph
        para = _Paragraph(target_p, text_frame)
        replace_paragraph_text(para, line)


def set_title(text_box, title_text):
    """Replace the title text box content while preserving its run formatting."""
    if not text_box.has_text_frame:
        return
    tf = text_box.text_frame
    # Take the first paragraph; squash everything into one paragraph.
    paragraphs = tf.paragraphs
    if not paragraphs:
        return
    first = paragraphs[0]
    body_elt = first._p.getparent()
    for p in paragraphs[1:]:
        body_elt.remove(p._p)
    replace_paragraph_text(first, title_text)


def normalize_run_spacing(text_frame):
    """Clear any letter-spacing (a:rPr/spc) on every run so text doesn't
    render with the template's stretched tracking. Also pin font name to
    Arial which the template's master prescribes."""
    from pptx.oxml.ns import qn
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            rPr = run._r.find(qn('a:rPr'))
            if rPr is not None and rPr.get('spc'):
                del rPr.attrib['spc']


def set_two_col_left(text_box, lines):
    """The two-column slides have a left text box with two paragraphs of body
    text and a 'SUBHEADING' header in the original. We squash all to one block
    of body paragraphs (no SUBHEADING heading).
    """
    if not text_box.has_text_frame:
        return
    tf = text_box.text_frame
    # The text frame in the template has structure:
    #   <p>SUBHEADING</p>  <p>body paragraph 1</p>  <p>body paragraph 2</p>
    # We want to replace with N body paragraphs, all using the body run's
    # formatting. Find the body paragraph run formatting (paragraph 1, idx 1
    # if 0 is SUBHEADING; otherwise paragraph 0).
    paragraphs = tf.paragraphs
    if not paragraphs:
        return
    # Pick a representative body paragraph: the longest one usually has run
    # formatting we want.
    rep_idx = 0
    rep_len = 0
    for i, p in enumerate(paragraphs):
        if len(p.text) > rep_len and p.runs:
            rep_len = len(p.text)
            rep_idx = i
    rep_xml = copy.deepcopy(paragraphs[rep_idx]._p)
    # Clear all paragraphs and rebuild from the rep template.
    body_elt = paragraphs[0]._p.getparent()
    for p in list(paragraphs):
        body_elt.remove(p._p)
    # Append N clones, set each line.
    from pptx.text.text import _Paragraph
    for line in lines:
        new_p = copy.deepcopy(rep_xml)
        body_elt.append(new_p)
        para = _Paragraph(new_p, tf)
        replace_paragraph_text(para, line)
    normalize_run_spacing(tf)


def insert_image_in_placeholder(slide, placeholder_shape, image_path):
    """Replace the placeholder rectangle with an inserted image at the same
    bounding box, preserving aspect ratio (centered)."""
    from PIL import Image

    img = Image.open(image_path)
    img_w, img_h = img.size
    img_aspect = img_w / img_h

    box_left = placeholder_shape.left
    box_top = placeholder_shape.top
    box_w = placeholder_shape.width
    box_h = placeholder_shape.height
    box_aspect = box_w / box_h

    if img_aspect > box_aspect:
        # Image is wider; fit width, shrink height
        new_w = box_w
        new_h = int(box_w / img_aspect)
        new_left = box_left
        new_top = box_top + (box_h - new_h) // 2
    else:
        # Image is taller; fit height, shrink width
        new_h = box_h
        new_w = int(box_h * img_aspect)
        new_top = box_top
        new_left = box_left + (box_w - new_w) // 2

    # Remove the placeholder shape itself (its cream rectangle).
    sp = placeholder_shape._element
    sp.getparent().remove(sp)
    # Insert the image at the computed bbox.
    slide.shapes.add_picture(image_path, new_left, new_top, new_w, new_h)


def remove_shape(shape):
    sp = shape._element
    sp.getparent().remove(sp)


def set_notes(slide, text):
    if not text:
        return
    notes_tf = slide.notes_slide.notes_text_frame
    notes_tf.text = text


# === Build =================================================================

prs = Presentation(SKELETON)
slides = list(prs.slides)
assert len(slides) == 14, f"Expected 14 slides, got {len(slides)}"

for i, (slide, content) in enumerate(zip(slides, SLIDES), start=1):
    kind = content["kind"]
    print(f"Slide {i}: kind={kind}")

    if kind == "title":
        # Slide 1 shapes:
        #   [0] 'Presentation' (subtitle line of title)
        #   [1] 'Title of'
        #   [2] 'DATE HERE'
        #   [3] 'PRESENTER NAME'
        #   [4] GROUP - footer with logo + 'College Name Goes Here'
        from pptx.util import Pt
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            txt = shape.text_frame.text
            if "Title of" in txt:
                set_title(shape, content["title_top"])
            elif "Presentation" == txt.strip():
                # Subtitle: handle multi-line and shrink font so it fits the
                # box without overlapping the title above.
                tf = shape.text_frame
                tf.word_wrap = True
                # Clear all existing paragraphs.
                paragraphs = tf.paragraphs
                template_xml = copy.deepcopy(paragraphs[0]._p)
                body_elt = paragraphs[0]._p.getparent()
                for p in paragraphs:
                    body_elt.remove(p._p)
                from pptx.text.text import _Paragraph
                for i, line in enumerate(content["title_main"].split("\n")):
                    new_p = copy.deepcopy(template_xml)
                    body_elt.append(new_p)
                    para = _Paragraph(new_p, tf)
                    replace_paragraph_text(para, line)
                    # Shrink the font so two-line subtitle fits inside the
                    # original 1.36-in-tall box.
                    for run in para.runs:
                        run.font.size = Pt(36)
            elif "DATE HERE" in txt:
                set_title(shape, content["date"])
            elif "PRESENTER NAME" in txt:
                set_title(shape, content["presenters"])
        # The "TEXAS A&M UNIVERSITY / College Name Goes Here" line in the
        # bottom-right corner is a baked PNG inside the GROUP shape -- the
        # text is rasterized, so we cannot edit it in-place. Cover it with a
        # white text-on-cream box that hides the placeholder while keeping
        # the maroon TAMU logo to its left.
        from pptx.util import Pt, Emu
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
        from pptx.enum.shapes import MSO_SHAPE
        # The image occupies roughly x=9.94 to x=12.27 in, y=6.21 to y=7.30
        # in. Cover only the text portion (right of the logo).
        cover_left = Emu(int(9.95 * 914400))
        cover_top = Emu(int(6.30 * 914400))
        cover_w = Emu(int(2.30 * 914400))
        cover_h = Emu(int(0.95 * 914400))
        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                      cover_left, cover_top, cover_w, cover_h)
        rect.fill.solid()
        rect.fill.fore_color.rgb = RGBColor(0xFA, 0xF7, 0xF3)  # cream, matches template image
        rect.line.fill.background()
        # Text box on top.
        tb = slide.shapes.add_textbox(cover_left, cover_top, cover_w, cover_h)
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.10)
        tf.margin_right = Inches(0.10)
        tf.margin_top = Inches(0.05)
        tf.margin_bottom = Inches(0.05)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        for i, line in enumerate(content["footer"].split("\n")):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            run = p.add_run()
            run.text = line
            run.font.name = "Arial"
            run.font.bold = (i == 0)
            run.font.size = Pt(12)
            run.font.color.rgb = RGBColor(0x50, 0x00, 0x00)
            p.alignment = PP_ALIGN.CENTER

    elif kind == "two_col":
        # Shapes: [0] picture placeholder, [1] left text, [2] title
        title_shape = None
        body_shape = None
        pic_shape = None
        for shape in slide.shapes:
            txt = shape.text_frame.text if shape.has_text_frame else ""
            if "Slide Title" in txt:
                title_shape = shape
            elif "Aliquam sollicitudin" in txt or "Lorem" in txt or "Vivamus" in txt:
                body_shape = shape
            elif shape.shape_type == 14:  # PLACEHOLDER (picture)
                pic_shape = shape
        if title_shape:
            set_title(title_shape, content["title"])
        if body_shape:
            set_two_col_left(body_shape, content["left"])
            # Optional per-slide font-size override to keep long bodies in box.
            override_pt = content.get("left_font_pt")
            if override_pt:
                from pptx.util import Pt
                for paragraph in body_shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(override_pt)
        # Image
        if pic_shape:
            if content.get("image"):
                insert_image_in_placeholder(slide, pic_shape, content["image"])
            elif content.get("image_caption"):
                # Replace placeholder with a centered stat-callout text block
                # using maroon text on the existing cream rectangle (so the
                # callout reads instead of being cream-on-cream).
                from pptx.util import Pt
                from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
                from pptx.enum.shapes import MSO_SHAPE
                left = pic_shape.left
                top = pic_shape.top
                w = pic_shape.width
                h = pic_shape.height
                remove_shape(pic_shape)
                # Draw a cream rectangle to mirror the template's image
                # placeholder background.
                rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
                rect.fill.solid()
                rect.fill.fore_color.rgb = RGBColor(0xFA, 0xF7, 0xF3)  # template cream
                rect.line.fill.background()  # no border
                # Text box on top with maroon text.
                tb = slide.shapes.add_textbox(left, top, w, h)
                tf = tb.text_frame
                tf.word_wrap = True
                tf.margin_left = Inches(0.4)
                tf.margin_right = Inches(0.4)
                tf.margin_top = Inches(0.5)
                tf.margin_bottom = Inches(0.5)
                tf.vertical_anchor = MSO_ANCHOR.MIDDLE
                lines = content["image_caption"].split("\n")
                for i, line in enumerate(lines):
                    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                    run = p.add_run()
                    run.text = line
                    run.font.name = "Arial"
                    if line.startswith("$"):
                        run.font.bold = True
                        run.font.size = Pt(56)
                    else:
                        run.font.bold = False
                        run.font.size = Pt(20)
                    run.font.color.rgb = RGBColor(0x50, 0x00, 0x00)  # TAMU maroon
                    p.alignment = PP_ALIGN.CENTER
            else:
                remove_shape(pic_shape)

    elif kind == "three_col":
        # Shapes: [0,1,2] picture placeholders (bottom), [3] title,
        # [4,5,6] subheading text columns
        title_shape = None
        col_text_shapes = []
        pic_shapes = []
        for shape in slide.shapes:
            txt = shape.text_frame.text if shape.has_text_frame else ""
            if "Slide Title" in txt:
                title_shape = shape
            elif "SUBHEADING" in txt:
                col_text_shapes.append(shape)
            elif shape.shape_type == 14:
                pic_shapes.append(shape)
        # Sort columns by x position so col 0 is leftmost
        col_text_shapes.sort(key=lambda s: s.left)
        if title_shape:
            set_title(title_shape, content["title"])
        for j, col in enumerate(content["cols"]):
            if j < len(col_text_shapes):
                tb = col_text_shapes[j]
                tf = tb.text_frame
                paragraphs = tf.paragraphs
                if len(paragraphs) < 2:
                    continue
                # First paragraph = column head (preserve template's bold style)
                replace_paragraph_text(paragraphs[0], col["head"].upper())
                # Second paragraph = first body line. Body can be a string or a list.
                body = col["body"]
                if isinstance(body, str):
                    body_lines = [body]
                else:
                    body_lines = list(body)
                # Drop any paragraphs after the second so we have a clean slate.
                body_elt = paragraphs[0]._p.getparent()
                for p in paragraphs[2:]:
                    body_elt.remove(p._p)
                # Set first body line in the existing body paragraph.
                replace_paragraph_text(paragraphs[1], body_lines[0])
                # Append additional body paragraphs as clones of paragraph[1].
                if len(body_lines) > 1:
                    template_xml = copy.deepcopy(paragraphs[1]._p)
                    from pptx.text.text import _Paragraph
                    for line in body_lines[1:]:
                        new_p = copy.deepcopy(template_xml)
                        body_elt.append(new_p)
                        para = _Paragraph(new_p, tf)
                        replace_paragraph_text(para, line)
                normalize_run_spacing(tf)
        # Remove the bottom picture placeholders -- we don't have icons.
        for pic in pic_shapes:
            remove_shape(pic)

    elif kind == "section":
        # Slide 11: Section Header layout, single text box "BREAKER TITLE PAGE"
        for shape in slide.shapes:
            if shape.has_text_frame and "BREAKER" in shape.text_frame.text:
                set_title(shape, content["title"])

    elif kind == "bullets":
        # Slide 12: Title and Content layout. Shape [0] is body with 3
        # 'OBJECTIVE Lorem...' paragraphs; shape [1] is title.
        title_shape = None
        body_shape = None
        for shape in slide.shapes:
            txt = shape.text_frame.text if shape.has_text_frame else ""
            if "Slide Title" in txt:
                title_shape = shape
            elif "OBJECTIVE" in txt or "Lorem" in txt:
                body_shape = shape
        if title_shape:
            set_title(title_shape, content["title"])
        if body_shape:
            set_two_col_left(body_shape, content["bullets"])

    elif kind == "thanks":
        # Slide 14: 'Texas A&M University / Division of Marketing & Communication'
        for shape in slide.shapes:
            if shape.has_text_frame:
                set_title(shape, content["footer"])

    set_notes(slide, content.get("notes", ""))

OUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(str(OUT))
print(f"\nWrote {OUT}")
