"""
Build the 48x36 TAMU Engineering showcase poster for STAR.

Corrected to the real hardware/software stack: RPi5 + RX72N + RPLiDAR C1 +
IMX219-83 stereo + two IMUs + HC-SR04 + slam_toolbox async + Nav2 +
m-explore-ros2 frontier exploration + SPI 10 Mbps with HARQ/FEC + nanopb.
Compliance engine scope is honest: one check implemented (ramp slope),
two stretch (trip, path width), four architected for deployment phase.

Usage:
    cd final_capstone
    bash charts/generate_all.sh
    python3 poster/build_poster.py
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

HERE = Path(__file__).resolve().parent
ROOT = HERE.parent
TEMPLATE = ROOT / "48x36_Poster_Template.pptx"
CHARTS = ROOT / "charts" / "output"
OUT = HERE / "STAR_Poster_48x36.pptx"

MAROON = RGBColor(0x50, 0x00, 0x00)
CHARCOAL = RGBColor(0x1a, 0x1a, 0x2e)
BLACK = RGBColor(0x11, 0x11, 0x11)
GREEN = RGBColor(0x2a, 0x9d, 0x8f)
GOLD = RGBColor(0xc2, 0x6e, 0x1a)


def set_text(tf, blocks):
    tf.clear()
    for i, b in enumerate(blocks):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = b.get("align", PP_ALIGN.LEFT)
        run = p.add_run()
        run.text = b["text"]
        run.font.size = Pt(b.get("size", 18))
        run.font.bold = b.get("bold", False)
        run.font.color.rgb = b.get("color", BLACK)
        if b.get("space_after"):
            p.space_after = Pt(b["space_after"])


def main():
    prs = Presentation(str(TEMPLATE))
    slide = prs.slides[0]

    shapes = {s.name: s for s in slide.shapes}
    title = shapes["Title 1"]
    col1 = shapes["Text Placeholder 12"]
    col2 = shapes["Text Placeholder 23"]
    col3 = shapes["Text Placeholder 24"]
    col4 = shapes["Text Placeholder 25"]

    # ---- Title ----
    set_text(title.text_frame, [
        {"text": "STAR: Spatial Topography Accessibility Robot",
         "bold": True, "size": 88, "color": MAROON, "align": PP_ALIGN.CENTER},
        {"text": "", "size": 14},
        {"text": "A distributed autonomous ground-robotics platform for continuous indoor ADA compliance measurement",
         "size": 34, "color": CHARCOAL, "align": PP_ALIGN.CENTER},
        {"text": "", "size": 14},
        {"text": "Team Locked Inc.  |  TAMU ESET Senior Capstone, Spring 2026",
         "bold": True, "size": 26, "align": PP_ALIGN.CENTER},
        {"text": "Department of Engineering Technology & Industrial Distribution",
         "size": 22, "align": PP_ALIGN.CENTER},
    ])

    # ---- Column 1: Problem Definition + Methodology ----
    set_text(col1.text_frame, [
        {"text": "Problem Definition", "bold": True, "size": 38, "color": MAROON},
        {"text": "", "size": 8},
        {"text": "28.7% of U.S. adults (~70M) live with a disability (CDC 2022 BRFSS).", "size": 20},
        {"text": "Industry estimate: ~73% of U.S. commercial buildings fail >=1 ADA standard.", "size": 20},
        {"text": "8,800 ADA Title III federal lawsuits filed in 2024 (+7% YoY, Seyfarth Shaw).", "size": 20},
        {"text": "DOJ penalties: $118,225 first / $236,451 subsequent (90 FR 29445, 2025).", "size": 20, "bold": True},
        {"text": "Only 32% of buyers ask about ADA at commercial inspection (Focus Bldg. Insp.).", "size": 20},
        {"text": "", "size": 10},
        {"text": "The bottleneck is measurement: a CASp audit costs $800-$50,000+ per building and takes 2-5 days.",
         "size": 20},
        {"text": "TAMU College Station: 5,200 acres, hundreds of buildings. Auditing every one manually is infeasible.",
         "size": 20, "bold": True},
        {"text": "", "size": 14},
        {"text": "Methodology", "bold": True, "size": 38, "color": MAROON},
        {"text": "", "size": 8},
        {"text": "Distributed platform: Raspberry Pi 5 high-level + custom RX72N PCB real-time, connected by SPI 10 Mbps with HARQ Chase Combining + FEC (K=7, rate-1/2) + nanopb Protocol Buffers over a CRC-32 frame protocol.",
         "size": 18},
        {"text": "", "size": 8},
        {"text": "Sensors (all physically wired):", "bold": True, "size": 22},
        {"text": "- SLAMTEC RPLiDAR C1 (360, 10 Hz, 12 m, +/- 3 cm)", "size": 18},
        {"text": "- Waveshare IMX219-83 stereo (dual Sony IMX219 8 MP, 60 mm baseline)", "size": 18},
        {"text": "- BNO055 9-DoF IMU (on RX72N) + ICM20948 9-DoF IMU (on camera)", "size": 18},
        {"text": "- 4 x HC-SR04 ultrasonic, DS18B20, BMP280", "size": 18},
        {"text": "", "size": 8},
        {"text": "Motor control: 4 x DFRobot FIT0520 (6 V, 210 RPM, 341 PPR Hall) + 4 x TI DRV8263H H-bridges + GPTW 32-bit PWM @ 20 kHz. RX72N runs 250 Hz discrete-time PID; Pi5 issues 100 Hz velocity commands.",
         "size": 18},
        {"text": "", "size": 8},
        {"text": "ROS2 Jazzy: slam_toolbox async + robot_localization EKF + Nav2 + m-explore-ros2 frontier exploration. 143 tests, 0 failures across 4 packages.",
         "size": 18},
    ])

    # ---- Column 2: Engineering Analysis (Seven ADA checks + responsibilities) ----
    set_text(col2.text_frame, [
        {"text": "Engineering Analysis", "bold": True, "size": 38, "color": MAROON},
        {"text": "", "size": 8},
        {"text": "Seven ADA 2010 geometric checks, implementation status labeled:", "size": 20, "bold": True},
        {"text": "", "size": 10},
        {"text": "[IMPLEMENTED]  1. Ramp slope > 1:12 (4.76 deg) - ADA 405.2", "size": 18, "color": GREEN, "bold": True},
        {"text": "   Open3D RANSAC plane normal + BNO055 pitch cross-validation (+/- 0.5 deg gate)", "size": 15},
        {"text": "", "size": 6},
        {"text": "[STRETCH]  2. Trip hazard > 0.25 inch - ADA 303", "size": 18, "color": GOLD, "bold": True},
        {"text": "   LiDAR vertical delta + BNO055 jolt correlation during traversal", "size": 15},
        {"text": "", "size": 6},
        {"text": "[STRETCH]  3. Accessible path width < 36 inches - ADA 403.5", "size": 18, "color": GOLD, "bold": True},
        {"text": "   Medial-axis transform on slam_toolbox occupancy grid", "size": 15},
        {"text": "", "size": 6},
        {"text": "[ARCHITECTED]  4. Ramp width < 36 inches - ADA 405.5", "size": 18},
        {"text": "   Plane-polygon bounding + principal-axis projection", "size": 15},
        {"text": "[ARCHITECTED]  5. Ramp landing < 60 x 60 inches - ADA 405.7", "size": 18},
        {"text": "   Inscribed-rectangle test on flat adjacent regions", "size": 15},
        {"text": "[ARCHITECTED]  6. Door clear width < 32 inches - ADA 404.2.3", "size": 18},
        {"text": "   IMX219-83 stereo SGBM + door-frame detection at handle height", "size": 15},
        {"text": "[ARCHITECTED]  7. Door threshold > 0.5 inch - ADA 404.2.5", "size": 18},
        {"text": "   LiDAR + stereo vertical profile across door bottom edge", "size": 15},
        {"text": "", "size": 14},
        {"text": "Why two IMUs?", "bold": True, "size": 22, "color": MAROON},
        {"text": "BNO055 sits on the chassis and reads robot-body pitch (authoritative for ramp slope).",
         "size": 16},
        {"text": "ICM20948 rides on the stereo camera board and gives the visual pipeline its own inertial reference for future visual-inertial odometry.",
         "size": 16},
    ])

    # ---- Column 3: Outcomes ----
    set_text(col3.text_frame, [
        {"text": "Outcomes", "bold": True, "size": 38, "color": MAROON},
        {"text": "", "size": 8},
        {"text": "Platform (complete)", "bold": True, "size": 26, "color": MAROON},
        {"text": "143/143 ROS2 tests passing across 4 packages.", "size": 20},
        {"text": "slam_toolbox async produces /map at 0.5 Hz on Raspberry Pi 5.", "size": 20},
        {"text": "Full TF chain: map -> odom -> base_link -> laser_frame.", "size": 20},
        {"text": "Nav2 autonomous goal execution validated.", "size": 20},
        {"text": "m-explore-ros2 frontier exploration to completion (no pre-loaded map).", "size": 20},
        {"text": "SPI transport: 1.6% utilization at 10 Mbps; HARQ/FEC protects against bit errors.", "size": 20},
        {"text": "250 Hz discrete-time PID motor loop, tuned against first-order model (tau = 75 ms).", "size": 20},
        {"text": "", "size": 10},
        {"text": "Compliance engine (this capstone)", "bold": True, "size": 26, "color": MAROON},
        {"text": "Ramp slope check implemented end-to-end.", "size": 20, "bold": True, "color": GREEN},
        {"text": "Wixey WR300 digital angle gauge ground truth (+/- 0.1 deg).", "size": 20},
        {"text": "On-campus ramp, n measurements logged to validation_log.csv.", "size": 20},
        {"text": "BNO055 pitch vs. LiDAR plane normal agreement gate +/- 0.5 deg.", "size": 20},
        {"text": "Auto-generated PDF audit report with georeferenced flag + sensor evidence.", "size": 20},
        {"text": "", "size": 10},
        {"text": "Live numbers on the printed poster come from the dress-rehearsal validation log, not from estimates. See extras/STAR_ValidationLog.xlsx.",
         "size": 16, "color": CHARCOAL},
    ])

    # ---- Column 4: Impact + References + Acknowledgments ----
    set_text(col4.text_frame, [
        {"text": "Impact", "bold": True, "size": 38, "color": MAROON},
        {"text": "", "size": 8},
        {"text": "At TAMU's 5,200-acre campus, continuous ADA measurement at amortized hardware cost is the same economic shift cybersecurity scanners brought to IT compliance a decade ago.",
         "size": 20},
        {"text": "", "size": 8},
        {"text": "Stakeholders", "bold": True, "size": 24, "color": MAROON},
        {"text": "- Department of Disability Resources: user advocacy", "size": 18},
        {"text": "  471 Houston Street, Student Services Bldg, Suite 122", "size": 15},
        {"text": "- Office of Risk, Ethics & Compliance: compliance ownership", "size": 18},
        {"text": "  orec.tamu.edu/ada", "size": 15},
        {"text": "- Facilities Services / Office of the University Architect: remediation", "size": 18},
        {"text": "", "size": 10},
        {"text": "Open-source release: github.com/Locked-Inc/STAR", "size": 20, "bold": True},
        {"text": "", "size": 14},
        {"text": "References", "bold": True, "size": 26, "color": MAROON},
        {"text": "Seyfarth Shaw ADA Title III Report, March 2025.", "size": 14},
        {"text": "Federal Register 90 FR 29445 (July 2025) - DOJ civil penalty adjustment.", "size": 14},
        {"text": "CDC Disability and Health Data System, 2022 BRFSS (released July 2024).", "size": 14},
        {"text": "Okoro et al., MMWR 2018; 67(32):882-887.", "size": 14},
        {"text": "Saha et al., Project Sidewalk, CHI 2019 (Best Paper). ACM 10.1145/3290605.3300292.", "size": 14},
        {"text": "Weld et al., ASSETS 2019 (Best Student Paper). ACM 10.1145/3308561.3353798.", "size": 14},
        {"text": "Hara et al., Tohme, UIST 2014. ACM 10.1145/2642918.2647403.", "size": 14},
        {"text": "Sun et al., Sensors 2024, 24(21):6828 - Stereo+LiDAR Loosely Coupled SLAM.", "size": 14},
        {"text": "Lang et al., Gaussian-LIC, arXiv:2404.06926, ICRA 2025.", "size": 14},
        {"text": "Turkan & Che, Automated ADA Curb-Ramp Assessment, Oregon State PacTrans.", "size": 14},
        {"text": "ADA 2010 Standards for Accessible Design, sections 303, 403.5, 404.2, 405.", "size": 14},
        {"text": "Full BibTeX: final_capstone/extras/references.bib. Source matrix: research/source_verification.md.",
         "size": 13, "color": CHARCOAL},
        {"text": "", "size": 10},
        {"text": "Acknowledgments", "bold": True, "size": 24, "color": MAROON},
        {"text": "TAMU ESET capstone faculty, TAMU Department of Engineering Technology & Industrial Distribution, TAMU Disability Resources for early stakeholder input.",
         "size": 15},
    ])

    # ---- Embed charts along the bottom of each column ----
    donut = CHARTS / "chart_b_compliance_donut.png"
    sankey = CHARTS / "chart_f_pipeline_sankey.png"
    resp = CHARTS / "chart_e_sensor_responsibilities.png"
    trend = CHARTS / "chart_a_lawsuit_trend.png"

    if donut.exists():
        slide.shapes.add_picture(str(donut),
                                 left=Inches(1.01), top=Inches(28.0),
                                 width=Inches(10.75), height=Inches(5.5))
    if sankey.exists():
        slide.shapes.add_picture(str(sankey),
                                 left=Inches(12.77), top=Inches(28.0),
                                 width=Inches(10.75), height=Inches(5.5))
    if resp.exists():
        slide.shapes.add_picture(str(resp),
                                 left=Inches(24.53), top=Inches(28.0),
                                 width=Inches(10.75), height=Inches(5.5))
    if trend.exists():
        slide.shapes.add_picture(str(trend),
                                 left=Inches(36.25), top=Inches(28.0),
                                 width=Inches(10.75), height=Inches(5.5))

    prs.save(str(OUT))
    print(f"wrote {OUT}")


if __name__ == "__main__":
    main()
