"""
Build the 12-slide STAR capstone deck (16:9, Midnight Executive palette).

Reflects the real hardware/software stack (RPi5 + RX72N custom PCB + SPI
with HARQ/FEC + RPLiDAR C1 + IMX219-83 stereo + BNO055 + ICM20948 +
HC-SR04), the real SLAM stack (slam_toolbox + Nav2 + m-explore-ros2), and
an honest compliance-engine scope (ramp slope IMPLEMENTED; trip hazard +
path width STRETCH; four remaining ADA checks ARCHITECTED).

Usage:
    cd final_capstone
    bash charts/generate_all.sh
    python3 deck/build_deck.py

Output: deck/STAR_Deck.pptx
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

HERE = Path(__file__).resolve().parent
ROOT = HERE.parent
CHARTS = ROOT / "charts" / "output"
OUT = HERE / "STAR_Deck.pptx"

# Midnight Executive palette (matches chart theme)
BG = RGBColor(0x1a, 0x1a, 0x2e)
PANEL = RGBColor(0x16, 0x21, 0x3e)
TEXT = RGBColor(0xe8, 0xe8, 0xf0)
MUTED = RGBColor(0x90, 0x90, 0xa8)
ACCENT = RGBColor(0xe6, 0x39, 0x46)
GOLD = RGBColor(0xf4, 0xa2, 0x61)
GREEN = RGBColor(0x2a, 0x9d, 0x8f)
BLUE = RGBColor(0x4c, 0xc9, 0xf0)
MAROON = RGBColor(0x50, 0x00, 0x00)


def dark_slide(prs, title_text, subtitle_text=None, speaker=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    bg.line.fill.background()
    bg.shadow.inherit = False

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.9))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = title_text
    r.font.size = Pt(34)
    r.font.bold = True
    r.font.color.rgb = TEXT

    if subtitle_text:
        sb = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(12.33), Inches(0.5))
        sp = sb.text_frame.paragraphs[0]
        sr = sp.add_run()
        sr.text = subtitle_text
        sr.font.size = Pt(18)
        sr.font.color.rgb = MUTED

    if speaker:
        foot = slide.shapes.add_textbox(Inches(10.5), Inches(6.9), Inches(2.5), Inches(0.3))
        fp = foot.text_frame.paragraphs[0]
        fp.alignment = PP_ALIGN.RIGHT
        fr = fp.add_run()
        fr.text = speaker
        fr.font.size = Pt(11)
        fr.font.color.rgb = MUTED

    return slide


def add_bullets(slide, left_in, top_in, w_in, h_in, bullets, size=22, color=TEXT):
    tb = slide.shapes.add_textbox(Inches(left_in), Inches(top_in),
                                  Inches(w_in), Inches(h_in))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(8)
        r = p.add_run()
        r.text = b
        r.font.size = Pt(size)
        r.font.color.rgb = color
    return tb


def add_image(slide, path, left_in, top_in, w_in=None, h_in=None):
    kw = {}
    if w_in is not None:
        kw["width"] = Inches(w_in)
    if h_in is not None:
        kw["height"] = Inches(h_in)
    return slide.shapes.add_picture(str(path), Inches(left_in), Inches(top_in), **kw)


def big_stat(slide, text, left_in, top_in, size=120, color=ACCENT):
    tb = slide.shapes.add_textbox(Inches(left_in), Inches(top_in),
                                  Inches(8), Inches(2))
    p = tb.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = True
    r.font.color.rgb = color
    return tb


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # -------- Slide 1: Hook --------
    s = dark_slide(prs, "$236,451 per door",
                   "The federal penalty for a single ADA repeat violation (90 FR 29445, July 2025)",
                   speaker="Team Lead")
    big_stat(s, "8,800", 0.6, 2.2, size=160, color=ACCENT)
    add_bullets(s, 3.8, 2.6, 9.0, 4.0, [
        "ADA Title III federal lawsuits filed in 2024 (+7% YoY, Seyfarth Shaw)",
        "Max federal penalty per repeat violation: $236,451 ($118,225 first)",
        "28.7% of U.S. adults live with a disability (~70M, CDC 2022 BRFSS)",
        "",
        "The only way to find violations today: one person, a tape measure, a clipboard.",
        "We built the measurement layer for a better way.",
    ], size=20)

    # -------- Slide 2: The problem --------
    s = dark_slide(prs, "The problem: who even measures this?",
                   speaker="Team Lead")
    donut = CHARTS / "chart_b_compliance_donut.png"
    disab = CHARTS / "chart_d_disability_breakdown.png"
    if donut.exists():
        add_image(s, donut, 0.3, 1.5, w_in=6.3)
    if disab.exists():
        add_image(s, disab, 6.8, 1.5, w_in=6.3)
    add_bullets(s, 0.5, 6.2, 12.3, 1.0, [
        "Industry estimate: ~73% of U.S. commercial buildings fail >=1 ADA standard (Building Principles).  "
        "Only 32% of commercial property buyers ask (Focus Building Inspections).  "
        "CASp audit: $800-$50,000+ depending on scope, 2-5 days on-site.",
    ], size=14, color=MUTED)

    # -------- Slide 3: Why now --------
    s = dark_slide(prs, "Why now: enforcement is accelerating",
                   speaker="Member 2")
    trend = CHARTS / "chart_a_lawsuit_trend.png"
    if trend.exists():
        add_image(s, trend, 0.3, 1.4, w_in=8.5)
    add_bullets(s, 9.0, 1.6, 4.0, 5.5, [
        "8,800 filings in 2024",
        "CA: 3,252 (So Cal Equal",
        "  Access Group filed 2,598)",
        "NY: 2,220. FL: 1,627. TX: 224.",
        "",
        "DOJ Title II rule:",
        "WCAG 2.1 AA by April 2026",
        "for public entities,",
        "TAMU included.",
        "",
        "Physical access pressure follows.",
    ], size=15)

    # -------- Slide 4: Our solution (platform + compliance layer) --------
    s = dark_slide(prs, "STAR: a platform for continuous ADA measurement",
                   "Distributed robot + compliance engine, ROS2 Jazzy, open source",
                   speaker="Member 2")
    add_bullets(s, 0.5, 1.8, 6.3, 5.0, [
        "Platform (demo-ready today):",
        "- Raspberry Pi 5 + custom RX72N PCB",
        "- RPLiDAR C1 + IMX219-83 stereo + 2 IMUs",
        "- slam_toolbox async + Nav2 + frontier explore",
        "- SPI 10 Mbps + HARQ/FEC + nanopb",
        "- 143 ROS2 tests, 0 failures",
    ], size=20)
    add_bullets(s, 6.9, 1.8, 6.3, 5.0, [
        "Compliance engine (new for this capstone):",
        "- Ramp slope >1:12  [IMPLEMENTED]",
        "- Trip hazard >0.25\"  [STRETCH]",
        "- Path width <36\"  [STRETCH]",
        "- Ramp width / landing / door clear / threshold  [ARCHITECTED]",
        "- PDF audit report at scan complete",
    ], size=20)

    # -------- Slide 5: Tech stack --------
    s = dark_slide(prs, "Tech stack",
                   speaker="Hardware Lead")
    add_bullets(s, 0.5, 1.5, 6.3, 5.5, [
        "Hardware",
        "- Raspberry Pi 5 (ROS2 Jazzy + UI + compliance engine)",
        "- Custom RX72N PCB (ThreadX, 250 Hz PID, 4 MB Flash / 1 MB SRAM)",
        "- SLAMTEC RPLiDAR C1 (360, 10 Hz, 12 m)",
        "- Waveshare IMX219-83 (dual 8 MP, 60 mm baseline)",
        "- BNO055 IMU (on chassis) + ICM20948 IMU (on camera)",
        "- 4x HC-SR04 ultrasonic + DS18B20 + BMP280",
        "- 4x DFRobot FIT0520 motors + 4x TI DRV8263H drivers",
        "- GPTW 32-bit 20 kHz PWM, MTU/TPU quadrature encoders",
    ], size=16)
    add_bullets(s, 6.9, 1.5, 6.3, 5.5, [
        "Software",
        "- ROS2 Jazzy (star_bringup, star_spi_bridge, star_safety_monitor)",
        "- slam_toolbox async + robot_localization EKF",
        "- Nav2 (NavFn A* + DWB) + m-explore-ros2 frontier",
        "- Compliance engine: Python + Open3D + OpenCV + reportlab",
        "- Gateway: Go + gRPC-Web + WebSocket + HTTP",
        "- UI: React 19.2 + Vite + TypeScript + Zustand (20 panels)",
        "- Firmware: C23, ThreadX, nanopb, HARQ/FEC",
        "- Open source: github.com/Locked-Inc/STAR",
    ], size=16)

    # -------- Slide 6: System architecture --------
    s = dark_slide(prs, "System architecture (real pipeline)",
                   speaker="Hardware Lead")
    sankey = CHARTS / "chart_f_pipeline_sankey.png"
    if sankey.exists():
        add_image(s, sankey, 0.3, 1.4, w_in=12.7)

    # -------- Slide 7: Seven ADA checks --------
    s = dark_slide(prs, "Seven ADA checks: status vs. sensor",
                   speaker="Software Lead")
    fusion = CHARTS / "chart_e_sensor_responsibilities.png"
    if fusion.exists():
        add_image(s, fusion, 0.3, 1.4, w_in=12.7)

    # -------- Slide 8: Demo --------
    s = dark_slide(prs, "Demo: autonomous explore + live ramp-slope check",
                   "(or pre-recorded backup)",
                   speaker="Software Lead")
    add_bullets(s, 0.5, 1.8, 12.3, 5.0, [
        "Robot autonomously explores the on-campus hallway (no floor plan)",
        "slam_toolbox async builds the map on the projected RViz feed",
        "Nav2 handles cost-map inflation around the 32 cm robot envelope",
        "Ramp is encountered: LiDAR plane normal and BNO055 pitch agree within 0.5 deg",
        "Slope > 4.76 degrees -> violation flagged, pose logged, CSV row written",
        "Scan completes: compliance engine emits audit_report.pdf with the flag",
    ], size=20)

    # -------- Slide 9: Validation --------
    s = dark_slide(prs, "Validation: what we can defend",
                   speaker="Software Lead")
    add_bullets(s, 0.5, 1.5, 6.3, 5.5, [
        "Platform (complete):",
        "- 143 / 143 ROS2 tests passing",
        "- slam_toolbox map at 0.5 Hz on RPi5",
        "- Full TF: map -> odom -> base_link -> laser_frame",
        "- Nav2 autonomous goals executed",
        "- SPI frame protocol at 1.6% utilization",
        "",
        "Compliance (implemented one check):",
        "- Wixey WR300 angle-gauge ground truth",
        "- BNO055 vs. LiDAR-plane-normal agreement",
        "- On-campus ramp, n = 5 to 30 sessions",
        "- Numbers are what we measured, not what we hoped for",
    ], size=16)
    resp = CHARTS / "chart_e_sensor_responsibilities.png"
    if resp.exists():
        add_image(s, resp, 6.9, 1.6, w_in=6.3)

    # -------- Slide 10: Broader impact --------
    s = dark_slide(prs, "Broader impact",
                   speaker="Team Lead")
    add_bullets(s, 0.5, 1.5, 12.3, 5.5, [
        "TAMU College Station: 5,200 acres, hundreds of buildings.",
        "Any CASp audit pass is a six- or seven-figure cost and runs months of specialist time.",
        "STAR re-uses the same robot across buildings at amortized hardware cost.",
        "",
        "Stakeholder fit:",
        "- Department of Disability Resources (471 Houston St, Student Services Bldg, Suite 122) - user advocacy",
        "- Office of Risk, Ethics & Compliance (orec.tamu.edu/ada) - compliance program owner",
        "- Facilities Services / Office of the University Architect - remediation execution",
        "",
        "Code released open source: github.com/Locked-Inc/STAR",
        "Same economic shift cybersecurity scanners brought to IT compliance a decade ago.",
    ], size=16)

    # -------- Slide 11: Future work --------
    s = dark_slide(prs, "Future work (architected, not vaporware)",
                   speaker="Member 2")
    add_bullets(s, 0.5, 1.5, 12.3, 5.5, [
        "Complete the architected ADA checks (ramp width, landing, door clear, threshold)",
        "RTAB-Map full LiDAR + RGB-D fusion (in-progress on main)",
        "Multi-floor traversal via elevator integration",
        "ML fixture recognition (restrooms, fountains, signage) for ADAS Chapter 6",
        "Outdoor mode (GPS + RTK) for curb ramps and accessible parking",
        "Fleet management: campus-as-a-service deployment pilot with Facilities Services",
    ], size=20)

    # -------- Slide 12: Q&A --------
    s = dark_slide(prs, "Team Locked Inc.",
                   "github.com/Locked-Inc/STAR  -  Pilot building partners welcome",
                   speaker="All four")
    add_bullets(s, 0.5, 2.0, 12.3, 4.5, [
        "Team Lead      - research, system integration, presentation",
        "Member 2       - motor control, electrical, PCB",
        "Hardware Lead  - sensor integration, mechanical, validation",
        "Software Lead  - ROS2, compliance engine, SLAM",
        "",
        "Questions?",
    ], size=22)

    prs.save(str(OUT))
    print(f"wrote {OUT}")


if __name__ == "__main__":
    main()
